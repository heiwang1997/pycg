"""
Copyright 2022 by Jiahui Huang. All rights reserved.
This file is part of PyCG toolbox and is released under "MIT License Agreement".
Please see the LICENSE file that should have been included as part of this package.
"""

import math
import textwrap
from typing import List, Union

import numpy as np
from PIL import Image
from pathlib import Path
from abc import ABC
import scipy.stats
from io import BytesIO


class ImageOperation(ABC):
    """ Lazy execution image operations """
    def apply(self, img: np.ndarray):
        pass


class CropOperation(ImageOperation):
    def __init__(self, x: int, y: int, w: int, h: int):
        self.x = x
        self.y = y
        self.w = w
        self.h = h

    def apply(self, img: np.ndarray):
        return img[
               max(0, self.y): min(self.y + self.h, img.shape[0] - 1),
               max(0, self.x): min(self.x + self.w, img.shape[1] - 1)]

    def enlarge(self, pix: int):
        self.x -= pix
        self.y -= pix
        self.w += 2 * pix
        self.h += 2 * pix


def ensure_float_image(img: np.ndarray):
    if img.dtype == np.uint8:
        img = img.astype(float) / 255.
    if img.shape[2] == 3:
        alpha = np.ones((img.shape[0], img.shape[1], 1), dtype=float)
        img = np.concatenate([img, alpha], axis=-1)
    return img


def alpha_compositing(image_a: np.ndarray, image_b: np.ndarray):
    """
    Composite two images A over B with alpha channels.
    If alpha channel is not provided, then full opacity is assumed.
    See https://en.wikipedia.org/wiki/Alpha_compositing

    :param image_a: H x W x 3/4 float or uint8 array
    :param image_b: H x W x 3/4 float or uint8 array
    :return H x W x 4 float np array.
    """
    assert image_a.shape[0] == image_b.shape[0] and image_a.shape[1] == image_b.shape[1]
    image_a = ensure_float_image(image_a)
    image_b = ensure_float_image(image_b)
    alpha_a = image_a[:, :, 3:]
    alpha_b = image_b[:, :, 3:]
    color_a, color_b = image_a[:, :, :3], image_b[:, :, :3]
    alpha_comp = alpha_b * (1 - alpha_a)
    alpha_final = alpha_a + alpha_comp
    color_final = (color_a * alpha_a + color_b * alpha_comp) / (alpha_final + 1e-8)
    return np.concatenate([color_final, alpha_final], axis=2)


def gamma_transform(img: np.ndarray, gamma: float = 1.0, alpha_only: bool = False):
    """
    Apply gamma transform to an image x = x ** gamma.
    :param img: source image
    :param gamma: float, gamma value
    :param alpha_only: if True only change the alpha channel, used to shrink / enlarge shadow
    :return: numpy array
    """
    img = ensure_float_image(img)
    if alpha_only:
        assert img.ndim == 3 and img.shape[2] == 4, f"Cannot find alpha channel with size {img.shape}!"
        img = np.concatenate([img[..., :3], img[..., 3:4] ** gamma], axis=2)
    else:
        img = img ** gamma
    img = np.clip(img, a_min=0, a_max=1)
    return img


def auto_crop(img: np.ndarray, bound_color=None, tol: float = 2 / 255., enlarge: int = 0, return_op: bool = False):
    img = ensure_float_image(img)
    color_img = img[:, :, :3]

    if bound_color is None:
        # Auto determine boundary color
        bound_samples = np.concatenate([
            color_img[:2].reshape(-1, 3), color_img[-2:].reshape(-1, 3),
            color_img[:, :2].reshape(-1, 3), color_img[:, -2:].reshape(-1, 3)
        ], axis=0)
        bound_color = scipy.stats.mode(bound_samples, axis=0).mode[0]
    else:
        if isinstance(bound_color[0], int) or isinstance(bound_color[0], np.uint8):
            bound_color = [t / 255. for t in bound_color]
        bound_color = np.asarray(bound_color)[:3]

    dist_img = np.linalg.norm(color_img - bound_color[None, None, :], axis=-1)
    dist_img = dist_img < tol
    x_mask, y_mask = np.all(dist_img, axis=0), np.all(dist_img, axis=1)
    x_mask, y_mask = np.where(~x_mask)[0], np.where(~y_mask)[0]
    crop_op = CropOperation(x_mask[0], y_mask[0], x_mask[-1] - x_mask[0] + 1, y_mask[-1] - y_mask[0] + 1)
    crop_op.enlarge(enlarge)
    if return_op:
        return crop_op
    return crop_op.apply(img)


def place_image(child_img: np.ndarray, parent_img: np.ndarray, pos_x: int, pos_y: int):
    """
    Place a child image over a parent image, the position of the left-upper corner is (pos_x, pos_y)
    :param child_img:
    :param parent_img:
    :param pos_x:
    :param pos_y:
    :return:
    """
    child_img = ensure_float_image(child_img)
    parent_img = ensure_float_image(parent_img)

    child_w, child_h = child_img.shape[1], child_img.shape[0]
    parent_w, parent_h = parent_img.shape[1], parent_img.shape[0]

    if pos_x >= parent_w or pos_y >= parent_h:
        return parent_img

    if pos_x + child_w <= 0 or pos_y + child_h <= 0:
        return parent_img

    parent_start_x = max(0, pos_x)
    parent_end_x = min(pos_x + child_w, parent_w)
    child_start_x = max(0, -pos_x)
    child_end_x = min(child_w, parent_w - pos_x)

    parent_start_y = max(0, pos_y)
    parent_end_y = min(pos_y + child_h, parent_h)
    child_start_y = max(0, -pos_y)
    child_end_y = min(child_h, parent_h - pos_y)

    parent_img[parent_start_y:parent_end_y, parent_start_x:parent_end_x] = \
        alpha_compositing(child_img[child_start_y:child_end_y, child_start_x:child_end_x],
                          parent_img[parent_start_y:parent_end_y, parent_start_x:parent_end_x])

    return parent_img


def chessboard(img_w: int, img_h: int, square_len: int = 20, seam_len: int = 0, color_a=None, color_b=None, color_seam=None):
    if color_a is None:
        color_a = (0.5, 0.5, 0.5, 1.0)
    if color_b is None:
        color_b = (1.0, 1.0, 1.0, 1.0)
    if color_seam is None:
        color_seam = (1.0, 1.0, 1.0, 1.0)

    block_a = np.full((square_len, square_len, 4), color_a)
    block_b = np.full((square_len, square_len, 4), color_b)
    basic_block = np.concatenate([np.concatenate([block_a, block_b], axis=0),
                            np.concatenate([block_b, block_a], axis=0)], axis=1)

    if seam_len > 0:
        border_len = seam_len // 2
        basic_block[:, :border_len] = color_seam
        basic_block[:, -border_len:] = color_seam
        basic_block[:border_len, :] = color_seam
        basic_block[-border_len:, :] = color_seam
        basic_block[:, square_len-border_len:square_len+border_len] = color_seam
        basic_block[square_len-border_len:square_len+border_len, :] = color_seam

    basic_block = np.tile(basic_block, (math.ceil(img_h / basic_block.shape[0]),
                                        math.ceil(img_w / basic_block.shape[1]), 1))

    return basic_block[:img_h, :img_w]


def solid(img_w: int, img_h: int, color=None):
    if color is None:
        color = (1.0, 1.0, 1.0, 1.0)
    if len(color) == 3:
        color = list(color) + [1.0]
    return np.full((img_h, img_w, 4), color)


def vlayout_images(image_list: list, slot_heights: list = None, width: int = -1, gap: int = 0, halignment: str = 'center',
                   background: list = None, margin: int = 0):
    if background is None:
        background = [1.0, 1.0, 1.0, 0.0]

    if slot_heights is None:
        slot_heights = [-1 for _ in range(len(image_list))]

    assert len(image_list) == len(slot_heights), "Image and slot sizes must be of same size!"
    assert halignment in ['left', 'center', 'right']

    image_list = [ensure_float_image(t) for t in image_list]
    for si in range(len(slot_heights)):
        if slot_heights[si] == -1:
            slot_heights[si] = image_list[si].shape[0]
    if width == -1:
        width = max([t.shape[1] for t in image_list])

    canvas_height = sum(slot_heights) + (len(image_list) - 1) * gap + 2 * margin
    canvas = np.zeros((canvas_height, width, 4)) + np.asarray(background)[np.newaxis, np.newaxis, :]

    place_center_y = slot_heights[0] / 2. + margin
    for cur_img, cur_slot_h, next_slot_h in zip(image_list, slot_heights, slot_heights[1:] + [-1]):
        place_ly = place_center_y - cur_img.shape[0] / 2.
        if halignment == 'left':
            place_lx = 0
        elif halignment == 'center':
            place_lx = width / 2. - cur_img.shape[1] / 2.
        else:
            place_lx = width - cur_img.shape[1]
        canvas = place_image(cur_img, canvas, int(place_lx), int(place_ly))
        place_center_y += (cur_slot_h / 2.) + gap + (next_slot_h / 2.)

    return canvas


def hlayout_images(image_list: list, slot_widths: list = None, height: int = -1, gap: int = 0, valignment: str = 'center',
                   background: list = None, margin: int = 0):
    """
    Note that this function will not scale images. It should be done outside!
    :param image_list:
    :param slot_widths: if -1, then auto infer
    :param height: if -1, then auto infer
    :param gap:
    :param valignment:
    :param background: [RGBA] canvas background.
    :return:
    """
    if background is None:
        background = [1.0, 1.0, 1.0, 0.0]

    if slot_widths is None:
        slot_widths = [-1 for _ in range(len(image_list))]

    assert len(image_list) == len(slot_widths), "Image and slot sizes must be of same size!"
    assert valignment in ['top', 'center', 'bottom']

    image_list = [ensure_float_image(t) for t in image_list]
    for si in range(len(slot_widths)):
        if slot_widths[si] == -1:
            slot_widths[si] = image_list[si].shape[1]
    if height == -1:
        height = max([t.shape[0] for t in image_list])

    canvas_width = sum(slot_widths) + (len(image_list) - 1) * gap + 2 * margin
    canvas = np.zeros((height, canvas_width, 4)) + np.asarray(background)[np.newaxis, np.newaxis, :]
    # print(canvas.max())

    place_center_x = slot_widths[0] / 2. + margin
    for cur_img, cur_slot_w, next_slot_w in zip(image_list, slot_widths, slot_widths[1:] + [-1]):
        place_lx = place_center_x - cur_img.shape[1] / 2.
        if valignment == 'top':
            place_ly = 0
        elif valignment == 'center':
            place_ly = height / 2. - cur_img.shape[0] / 2.
        else:
            place_ly = height - cur_img.shape[0]
        canvas = place_image(cur_img, canvas, int(place_lx), int(place_ly))
        place_center_x += (cur_slot_w / 2.) + gap + (next_slot_w / 2.)

    return canvas


def text(text, font='DejaVuSansMono.ttf', font_size=16, max_width=None):
    from PIL import ImageFont, ImageDraw

    font_obj = ImageFont.truetype(font, font_size)

    # Get character width and determine text-warps
    if max_width is not None:
        c_width = font_obj.getsize('ABC')[0] // 3
        max_characters = math.floor(max_width / c_width)
        text = textwrap.wrap(text, width=max_characters)

    if not isinstance(text, list):
        text = text.split('\n')

    # Determine real dimensions:
    font_w, font_h = 0, 0
    for line_text in text:
        font_dim = font_obj.getsize(line_text)
        font_h += font_dim[1]
        font_w = max(font_w, font_dim[0])

    img = Image.new('RGB', (font_w, font_h), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    line_h = 0
    for line_text in text:
        draw.text((0, line_h), line_text, font=font_obj, fill=(0, 0, 0))
        line_h += font_obj.getsize(line_text)[1]

    img = np.asarray(img)
    return img


def assembled_slides(pic_matrix: List[List[Union[Path, str, np.ndarray]]], direction: str = "vertical",
                     x_block_cm: float = 4.0, x_gap_cm: float = 0.5,
                     y_block_cm: float = 4.0, y_gap_cm: float = 0.5):
    """
    Do a slide with a picture matrix, apply '.save(xxx.pptx)' to the output of this method to save
        Note For future: if you want to use a assembled group, please run it multiple times, and adjust positions.
    :param pic_matrix: matrix of pictures, could either be numpy, or picture paths.
    :param direction: either vertical or horizontal
        Note that if horizontal, pic_matrix will be treated as transposed.
    :param x_block_cm: width of each element
        This will be treated as the maximum width if direction == 'horizontal'
    :param x_gap_cm: gap of the elements in x direction
    :param y_block_cm: height of each element
        This will be treated as the maximum height if direction == 'vertical'
    :param y_gap_cm: gap of the elements in y direction
    :return: pptx.Presentation (requires python-pptx)
    """
    for pic_row in pic_matrix:
        assert len(pic_matrix[0]) == len(pic_row)

    from pptx import Presentation
    from pptx.util import Cm

    x_block, x_gap = Cm(x_block_cm), Cm(x_gap_cm)
    y_block, y_gap = Cm(y_block_cm), Cm(y_gap_cm)

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    def np_to_file(img: np.ndarray):
        f_obj = BytesIO()
        write(img, f_obj)
        f_obj.seek(0)
        return f_obj

    def get_row_data(input_row):
        out_data = []
        for t in input_row:
            if isinstance(t, str) or isinstance(t, Path):
                out_data.append(read(t))
            else:
                out_data.append(t)
        out_data = [ensure_float_image(t) for t in out_data]
        return out_data

    if direction == "vertical":
        y_ptr = Cm(0)
        for pic_row in pic_matrix:
            row_data = get_row_data(pic_row)
            row_w, row_h = row_data[0].shape[1], row_data[0].shape[0]
            if row_h / row_w > y_block / x_block:
                ppt_size = {'height': y_block}
                x_incr = (x_block - (y_block / row_h * row_w)) / 2.
                y_ptr_new = y_ptr + y_block
            else:
                ppt_size = {'width': x_block}
                x_incr = 0
                y_ptr_new = y_ptr + (x_block / row_w * row_h)

            for row_idx, row_img in enumerate(row_data):
                left = row_idx * (x_gap + x_block) + x_incr
                _ = slide.shapes.add_picture(np_to_file(row_img), left, y_ptr, **ppt_size)

            y_ptr = y_ptr_new + y_gap

        prs.slide_width = (x_block + x_gap) * len(pic_matrix[0]) - x_gap
        prs.slide_height = int(y_ptr - y_gap)

    elif direction == "horizontal":
        x_ptr = Cm(0)
        for pic_row in pic_matrix:
            row_data = get_row_data(pic_row)
            row_w, row_h = row_data[0].shape[1], row_data[0].shape[0]
            if row_w / row_h > x_block / y_block:
                ppt_size = {'width': x_block}
                y_incr = (y_block - (x_block / row_w * row_h)) / 2.
                x_ptr_new = x_ptr + x_block
            else:
                ppt_size = {'height': y_block}
                y_incr = 0
                x_ptr_new = x_ptr + (y_block / row_h * row_w)

            for row_idx, row_img in enumerate(row_data):
                top = row_idx * (y_gap + y_block) + y_incr
                _ = slide.shapes.add_picture(np_to_file(row_img), x_ptr, top, **ppt_size)

            x_ptr = x_ptr_new + x_gap

        prs.slide_width = int(x_ptr - x_gap)
        prs.slide_height = (y_block + y_gap) * len(pic_matrix[0]) - y_gap

    else:
        raise NotImplementedError

    return prs


def from_mplot(fig, close: bool = False):
    from matplotlib.backends.backend_agg import FigureCanvasAgg as FigureCanvas
    import matplotlib.pyplot as plt

    canvas = FigureCanvas(fig)
    canvas.draw()
    image = np.fromstring(canvas.tostring_rgb(), dtype='uint8')
    image = image.reshape(canvas.get_width_height()[::-1] + (3,))

    if close:
        plt.close(fig)

    return image


def show(*imgs, reverse_rgb: bool = False, n_rows: int = 1, subfig_size: int = 3, shrink_batch_dim: bool = False):
    import matplotlib.pyplot as plt

    n_cols = math.ceil(len(imgs) / n_rows)
    plt.figure(figsize=(n_cols * subfig_size, n_rows * subfig_size))
    show_imgs = []
    for img in imgs:
        if hasattr(img, 'cuda'):
            if shrink_batch_dim:
                img = img[0]
            img = img.detach().cpu().numpy()

        if img.ndim > 2 and img.shape[2] == 4:
            img_back = chessboard(img.shape[1], img.shape[0])
            img = alpha_compositing(img, img_back)
        if img.ndim > 2 and img.shape[2] >= 3 and reverse_rgb:
            img = np.copy(img)
            img[:, :, :3] = img[:, :, [2, 1, 0]]
        show_imgs.append(img)

    for img_id, img in enumerate(show_imgs):
        plt.subplot(n_rows, n_cols, img_id + 1)
        plt.imshow(show_imgs[img_id])
    plt.show()


def add_alpha_color(img: np.ndarray, alpha_color, tol: float = 2 / 255.):
    img = ensure_float_image(img)
    if isinstance(alpha_color[0], int) or isinstance(alpha_color[0], np.uint8):
        alpha_color = [t / 255. for t in alpha_color]

    origin_alpha = img[:, :, 3]
    dist_img = np.linalg.norm(img[:, :, :3] - np.asarray(alpha_color)[np.newaxis, np.newaxis, :], axis=-1)
    dist_img = dist_img < tol
    new_alpha = 1. - dist_img.astype(float)
    new_alpha = np.minimum(origin_alpha, new_alpha)
    img[:, :, 3] = new_alpha
    return img


def read(path):
    path = Path(path)
    if path.suffix == ".hdr" or path.suffix == ".exr":
        import imageio
        # https://imageio.readthedocs.io/en/v2.4.1/format_hdr-fi.html#hdr-fi
        return np.asarray(imageio.v2.imread(path))
    with Image.open(path) as im:
        return np.asarray(im)


def write(img: np.ndarray, path_or_f):
    # Allow for swapped arguments.
    if isinstance(path_or_f, np.ndarray):
        img, path_or_f = path_or_f, img

    if img.dtype == float:
        img = (img * 255).astype(np.uint8)
    img = Image.fromarray(img)

    if isinstance(path_or_f, Path) or isinstance(path_or_f, str):
        Path(path_or_f).parent.mkdir(parents=True, exist_ok=True)
        img.save(path_or_f)
    elif callable(getattr(path_or_f, "seek")):
        img.save(path_or_f, format="png")
    else:
        raise NotImplementedError
